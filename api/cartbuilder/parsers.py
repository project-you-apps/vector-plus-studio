"""
parsers.py — Document text extraction for Cart Builder GUI

File-type coverage expanded 2026-06-14 to match the graphify (safishamsi/graphify, MIT)
file-type matrix Andy flagged 2026-05-11. Approach: use the same upstream Python
libraries graphify uses (BeautifulSoup, python-pptx, yaml, json) rather than
vendoring graphify itself. Code files get raw-text parsing for v1 — AST-aware
chunking via tree-sitter is a v2 improvement once chunking strategy is revisited
across the product suite.

Parser dependencies are imported lazily so that missing optional deps only break
their specific file type, not the whole module.
"""
import os
import re
import json as _json
from pathlib import Path


# Lone UTF-16 surrogates (U+D800–U+DFFF) are an encoding mechanism, not real
# characters. When they appear in Python strings they are decoder artifacts
# (e.g. a buggy UTF-16 path that advanced 2 bytes instead of 4 on an astral
# character, emitting the low surrogate as a stray adjacent "character" right
# after the properly-decoded codepoint). They crash FastAPI's JSON serializer
# at response time. Drop them — astral characters themselves (codepoints
# >= U+10000, like 🅿 U+1F17F) are OUTSIDE this range and are preserved.
_LONE_SURROGATE_RE = re.compile(r'[\ud800-\udfff]')


def scrub_lone_surrogates(text: str) -> str:
    return _LONE_SURROGATE_RE.sub('', text) if text else text


def _join_single_char_runs(text: str, min_run: int = 3) -> str:
    """Collapse runs of consecutive single-character lines into one line.

    Business forms set narrow column headers vertically, one glyph per line.
    PyMuPDF reads them top-to-bottom, so a Sysco invoice linearizes to:

        S
        P
        T
        L
        I
        LOC

    That is 9% of the extracted lines on a real invoice and it is pure layout
    artifact — it carries no meaning as separate lines, wastes embedding
    signal, and looks broken in the passage viewer. Joining the run to
    "S P T L I" keeps every character while costing five lines of noise.

    Requires min_run consecutive singles so ordinary text (a lone "I", "a",
    a list marker) is left alone.
    """
    lines = text.split("\n")
    out: list[str] = []
    run: list[str] = []

    def flush():
        if len(run) >= min_run:
            out.append(" ".join(run))
        else:
            out.extend(run)
        run.clear()

    for line in lines:
        if len(line.strip()) == 1:
            run.append(line.strip())
        else:
            flush()
            out.append(line)
    flush()
    return "\n".join(out)


def parse_pdf(filepath: Path) -> list[dict]:
    import fitz  # PyMuPDF
    doc = fitz.open(str(filepath))
    results = []
    for i, page in enumerate(doc):
        text = _join_single_char_runs(page.get_text().strip())
        if text:
            results.append({"text": text, "page": i + 1, "source": filepath.name})
    doc.close()
    return results


# Tables from text-layer PDFs.
#
# Until now, `tables` was only ever populated by the Image Builder OCR path,
# so a SCANNED document produced table patterns and a DIGITAL one did not.
# The better input produced the worse structure. `parse_pdf` above runs
# PyMuPDF get_text(), which linearizes the page — on a vendor invoice the
# headers come back as loose lines and every qty-to-description-to-price
# relationship is gone.
#
# Concretely, on a real Sysco invoice the OCR path returned a header reading
# `UNITTYTYTYTY...` and an extended price glued to a state fee as `86.61.58`,
# while the same invoice as a digital PDF yields exact rows. Invoices, order
# guides and statements are the core ingestion case and the table IS the
# content, so the text route needs to produce tables too.
#
# Emits the same {html, page, bbox} shape Docling returns, so everything
# downstream (`content_type: "table"` patterns, `_table_html_to_text`) is
# unchanged.

# Table extraction costs roughly 50-200ms/page, so long PDFs cap out. Matches
# the classify sampling cap for the same reason: predictable ingest latency.
PDF_TABLE_MAX_PAGES = 15
# Two cells in one row is a table; one column is a list, and layout scaffolding
# frequently parses as a single wide column.
PDF_TABLE_MIN_COLS = 2
PDF_TABLE_MIN_ROWS = 2


# Empty cells get a visible placeholder rather than being left blank. On a wide
# business form most cells are empty, and a markdown row of bare pipes gives a
# reader no way to count columns or tell which value sits under which header —
# the layout is the information. An interpunct is unobtrusive, unambiguously
# "nothing here", and keeps the grid legible. Aesthetics of the cell don't
# matter; the rows and columns reading correctly does.
EMPTY_CELL = "·"


def _cell_html(value) -> str:
    text = "" if value is None else str(value).replace("\n", " ").strip()
    if not text:
        return EMPTY_CELL
    return (text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;"))


_SPACED_CAPS = "letter-spaced header repair"


def _collapse_spaced(cell) -> str:
    """'L O C' -> 'LOC'. Business forms set vertical headers as spaced capitals.

    Only fires when EVERY token is a single character, so 'TOT.PCS', 'A 1' and
    ordinary prose are untouched. Without this, a Sysco invoice header reads
    'L O C | C O N T | T A X | P I' and no reader — human or model — can tell
    those are LOC / CONT / TAX / PI.
    """
    # _normalize_columns pads short rows with None, so coerce before splitting.
    text = "" if cell is None else str(cell)
    parts = text.split()
    if len(parts) > 1 and all(len(p) == 1 for p in parts):
        return "".join(parts)
    return text


def _numeric_dominant(cell) -> bool:
    """True for a VALUE cell, false for a sentence that happens to contain digits.

    '37.99', '6001208', '9.52OZ' -> True.
    'CALL 800-797-2627 OR EMAIL...' and the PACA legal paragraph -> False.
    Length cap plus a 50% digit ratio is what separates them; 'contains a
    digit' alone lets phone numbers and statute citations pose as data.
    """
    s = ("" if cell is None else str(cell)).strip()
    if not s or len(s) > 24:
        return False
    digits = sum(ch.isdigit() for ch in s)
    return bool(digits) and digits / len(s) >= 0.5


def _split_table_rows(rows: list[list]):
    """-> (header_row | None, data_rows, prose_rows).

    A ruled business form is one big bordered region, so pdfplumber returns the
    WHOLE PAGE as a single table: column headers, marketing copy, line items,
    totals, signature block and legal boilerplate all in one grid. Rendering
    that verbatim produces a 23-column wall where each prose sentence is
    smeared across eight cells.

    Split by shape rather than by keyword, so it generalizes past Sysco:
      - data row  : >= 2 numeric-dominant cells (real values in real columns)
      - header    : the first all-text row wide enough to be labels
      - prose     : everything else — emitted as TEXT, never dropped

    Nothing is discarded; prose is relocated out of the grid into the passage
    body where it reads normally.
    """
    header, data, prose = None, [], []
    for row in rows:
        filled = [c for c in row if c]
        if not filled:
            continue
        n_values = sum(1 for c in filled if _numeric_dominant(c))
        if header is None and n_values == 0 and len(filled) >= 4:
            header = row
            continue
        (data if n_values >= 2 else prose).append(row)
    return header, data, prose


def _merge_disjoint_columns(rows: list[list]) -> list[list]:
    """Merge adjacent columns never both filled in the same row.

    Lossless by construction: if two columns never co-occupy a row, one column
    can hold both without a collision. Business forms produce many such
    columns because different row types use different parts of the grid.
    Measured on real Sysco invoices: 23 columns -> 14, no content lost.
    """
    while True:
        width = max((len(r) for r in rows), default=0)
        merged = False
        for a in range(width - 1):
            b = a + 1
            if all(not (r[a] and r[b]) for r in rows):
                for r in rows:
                    r[a] = r[a] or r[b]
                    del r[b]
                merged = True
                break
        if not merged:
            return rows


def _normalize_columns(rows: list[list]) -> list[list]:
    """Pad ragged rows to a uniform width, then drop all-empty columns.

    The padding is the load-bearing half. pdfplumber returns ragged rows — on
    one real invoice a header row had 23 cells and a continuation row had 2 —
    and emitting those directly produces a malformed markdown table with
    misaligned columns.

    The all-empty-column drop is cheap insurance that measured ZERO hits
    across 52 tables from 99 real Sysco invoices: on those forms every column
    is populated in at least one row, so the table is genuinely 23 columns
    wide rather than padded. Kept for other document shapes, but do not
    expect it to tidy invoices — it won't, and the width you see is real.
    """
    width = max((len(r) for r in rows), default=0)
    padded = [list(r) + [None] * (width - len(r)) for r in rows]
    keep = [c for c in range(width)
            if any(row[c] is not None and str(row[c]).strip() for row in padded)]
    return [[row[c] for c in keep] for row in padded]


def extract_pdf_tables(filepath: Path) -> list[dict]:
    """Extract tables from a text-layer PDF -> [{html, page, bbox}].

    Never raises: a table-extraction failure must not fail an ingest that
    would otherwise succeed. On any error the caller gets [] and the document
    still lands as text, which is exactly today's behaviour.
    """
    try:
        import pdfplumber
    except ImportError:
        return []

    out: list[dict] = []
    try:
        with pdfplumber.open(str(filepath)) as pdf:
            for i, page in enumerate(pdf.pages[:PDF_TABLE_MAX_PAGES]):
                try:
                    found = page.find_tables()
                except Exception:
                    continue
                for t in found:
                    try:
                        rows = t.extract()
                    except Exception:
                        continue
                    rows = [r for r in (rows or [])
                            if any(c is not None and str(c).strip() for c in r)]
                    if len(rows) < PDF_TABLE_MIN_ROWS:
                        continue
                    if max((len(r) for r in rows), default=0) < PDF_TABLE_MIN_COLS:
                        continue
                    rows = _normalize_columns(rows)
                    if max((len(r) for r in rows), default=0) < PDF_TABLE_MIN_COLS:
                        continue

                    # Repair letter-spaced headers before anything reads them.
                    rows = [[_collapse_spaced(c) for c in r] for r in rows]

                    # Lift running prose out of the grid, then collapse the
                    # columns it was forcing open. Fall back to the whole grid
                    # whenever this doesn't clearly apply — a table that isn't
                    # a business form should come through exactly as before.
                    header, data_rows, prose_rows = _split_table_rows(rows)
                    notes: list[str] = []
                    if header is not None and len(data_rows) >= 1:
                        rows = _merge_disjoint_columns(
                            [r[:] for r in ([header] + data_rows)])
                        notes = [" ".join(c for c in r if c).strip()
                                 for r in prose_rows]
                        notes = [n for n in notes if n]

                    body = "".join(
                        "<tr>" + "".join(f"<td>{_cell_html(c)}</td>" for c in r) + "</tr>"
                        for r in rows)
                    out.append({
                        "html": f"<table>{body}</table>",
                        "page": i + 1,
                        "bbox": list(t.bbox) if getattr(t, "bbox", None) else [],
                        # Prose lifted out of the grid. The caller appends these
                        # below the rendered table so the passage keeps every
                        # word the page had, just not jammed into cells.
                        "notes": notes,
                    })
    except Exception:
        return []
    return out


# Day 2 — PDF classification for Image Builder routing. Mirrors the
# frontend classifyPdf() (frontend/src/cart-builder-v2/parsers/pdf.ts).
# Same constants both sides so paired and browser-only builds route the same
# file the same way (golden-path invariant).
PDF_CLASSIFY_TEXT_THRESHOLD = 500
# Sample up to 15 pages during classification. Some PDFs are clean on the
# first few pages and only expose broken ToUnicode fonts later on; head-only
# sampling misses that corruption. Cost is ~100-500ms extra classify time
# on a 15-page document; long PDFs cap here so classify stays fast.
PDF_CLASSIFY_MAX_PAGES = 15
# PDFs with broken font ToUnicode maps return LOTS of
# characters (well over the 500-char threshold) but most are Private Use
# Area / replacement / non-Latin garbage that produces unreadable ingest.
# Two-level check:
#   - Per-page: if ANY sampled page has substantial content (>50 chars)
#     but < PDF_CLASSIFY_PAGE_READABLE_THRESHOLD readable, route to
#     Image Builder. One bad page = broken font used elsewhere in the doc.
#   - Aggregate fallback: if overall readable fraction < threshold, same.
# 0.6 is the initial pick; tune with more samples if false-positives appear.
PDF_CLASSIFY_READABLE_THRESHOLD = 0.6
PDF_CLASSIFY_PAGE_READABLE_THRESHOLD = 0.6
PDF_CLASSIFY_PAGE_MIN_CHARS = 50


def _readable_char_count(text: str) -> int:
    """Count characters that look like real text: printable ASCII, common
    whitespace, or Latin-1/Extended-A/B codepoints (accents, etc). Everything
    else (PUA, replacement, most CJK) counts as unreadable for the purposes
    of ToUnicode-corruption detection."""
    count = 0
    for c in text:
        cp = ord(c)
        if 0x20 <= cp <= 0x7E:
            count += 1
        elif c in "\n\r\t":
            count += 1
        elif 0xA0 <= cp <= 0x24F:
            count += 1
    return count


def classify_pdf(filepath: Path) -> str:
    """Return 'text' if the PDF has extractable readable text, else 'scanned'.

    Sums PyMuPDF's page.get_text() lengths across the first
    PDF_CLASSIFY_MAX_PAGES pages, and independently sums the count of
    readable characters (printable ASCII + Latin extensions). Returns 'text'
    only when total > PDF_CLASSIFY_TEXT_THRESHOLD (500) AND the readable
    fraction >= PDF_CLASSIFY_READABLE_THRESHOLD (0.6). Otherwise routes to
    Image Builder /ocr.

    Failure mode: if PyMuPDF can't open the file, return 'scanned' — Docling
    is more likely to salvage a malformed PDF than PyMuPDF is, and if it
    can't either, the calling code surfaces an OCR failure the user can act
    on (rather than a silent parse-to-empty).
    """
    import fitz
    try:
        doc = fitz.open(str(filepath))
    except Exception:
        return "scanned"
    total_chars = 0
    readable_chars = 0
    corrupt_page_found = False
    try:
        pages_to_check = min(len(doc), PDF_CLASSIFY_MAX_PAGES)
        for i in range(pages_to_check):
            try:
                text = (doc[i].get_text() or "").strip()
                page_len = len(text)
                page_readable = _readable_char_count(text)
                total_chars += page_len
                readable_chars += page_readable
                # Per-page corruption check: substantial content but low
                # readable fraction signals a broken ToUnicode font used
                # somewhere in the document. One bad page taints the whole
                # ingest because the same font likely appears elsewhere.
                if page_len >= PDF_CLASSIFY_PAGE_MIN_CHARS:
                    page_fraction = page_readable / page_len
                    if page_fraction < PDF_CLASSIFY_PAGE_READABLE_THRESHOLD:
                        corrupt_page_found = True
            except Exception:
                # Skip page on error; partial score still useful — a truly
                # scanned PDF stays under threshold no matter which page
                # blows up.
                continue
    finally:
        try:
            doc.close()
        except Exception:
            pass
    readable_fraction = readable_chars / total_chars if total_chars else 0.0
    if total_chars <= PDF_CLASSIFY_TEXT_THRESHOLD:
        return "scanned"
    if corrupt_page_found:
        return "scanned"
    if readable_fraction < PDF_CLASSIFY_READABLE_THRESHOLD:
        return "scanned"
    return "text"


# Image-file extension set. Aligned with image-builder/main.py's
# SUPPORTED_FORMATS list. PDFs are handled via classify_pdf, not here.
_IMAGE_EXTENSIONS = {
    ".jpg", ".jpeg", ".png", ".heic", ".heif",
    ".tif", ".tiff", ".webp", ".bmp",
}


def is_image_file(filepath: Path) -> bool:
    """True when the file extension is one Image Builder can OCR directly.

    Used by builder.py's per-file routing to decide whether to POST straight
    to Image Builder /ocr (image) vs. classify a PDF first.
    """
    return filepath.suffix.lower() in _IMAGE_EXTENSIONS


def parse_docx(filepath: Path) -> list[dict]:
    import docx
    doc = docx.Document(str(filepath))
    text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    if not text.strip():
        return []
    return [{"text": text, "page": None, "source": filepath.name}]


def parse_xlsx(filepath: Path) -> list[dict]:
    import openpyxl
    wb = openpyxl.load_workbook(str(filepath), read_only=True, data_only=True)
    results = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            line = " | ".join(c for c in cells if c)
            if line.strip():
                rows.append(line)
        if rows:
            text = f"Sheet: {sheet_name}\n" + "\n".join(rows)
            results.append({"text": text, "page": None, "source": f"{filepath.name}:{sheet_name}"})
    wb.close()
    return results


def parse_markdown(filepath: Path) -> list[dict]:
    text = filepath.read_text(encoding="utf-8", errors="replace")
    sections = []
    current = []
    for line in text.split("\n"):
        if line.startswith("## ") and current:
            sections.append("\n".join(current))
            current = [line]
        else:
            current.append(line)
    if current:
        sections.append("\n".join(current))

    results = []
    for i, section in enumerate(sections):
        if section.strip():
            results.append({"text": section.strip(), "page": i + 1, "source": filepath.name})
    return results


def parse_text(filepath: Path) -> list[dict]:
    text = filepath.read_text(encoding="utf-8", errors="replace")
    if not text.strip():
        return []
    return [{"text": text.strip(), "page": None, "source": filepath.name}]


def parse_rtf(filepath: Path) -> list[dict]:
    from striprtf.striprtf import rtf_to_text
    raw = filepath.read_text(encoding="utf-8", errors="replace")
    text = rtf_to_text(raw)
    if not text.strip():
        return []
    return [{"text": text.strip(), "page": None, "source": filepath.name}]


def parse_html(filepath: Path) -> list[dict]:
    try:
        from bs4 import BeautifulSoup
    except ImportError:
        # Graceful fallback to raw text if BeautifulSoup not installed
        return parse_text(filepath)
    raw = filepath.read_text(encoding="utf-8", errors="replace")
    soup = BeautifulSoup(raw, "html.parser")
    for tag in soup(["script", "style", "noscript"]):
        tag.decompose()
    text = soup.get_text(separator="\n").strip()
    if not text:
        return []
    return [{"text": text, "page": None, "source": filepath.name}]


def parse_pptx(filepath: Path) -> list[dict]:
    try:
        from pptx import Presentation
    except ImportError:
        return []
    prs = Presentation(str(filepath))
    results = []
    for i, slide in enumerate(prs.slides):
        parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text and shape.text.strip():
                parts.append(shape.text.strip())
            if shape.has_table:
                for row in shape.table.rows:
                    line = " | ".join(c.text.strip() for c in row.cells if c.text.strip())
                    if line:
                        parts.append(line)
        if parts:
            text = "\n".join(parts)
            results.append({"text": text, "page": i + 1, "source": filepath.name})
    return results


def parse_yaml(filepath: Path) -> list[dict]:
    try:
        import yaml
    except ImportError:
        return parse_text(filepath)
    raw = filepath.read_text(encoding="utf-8", errors="replace")
    try:
        data = yaml.safe_load(raw)
    except Exception:
        return parse_text(filepath)
    text = _flatten_structured(data)
    if not text.strip():
        return []
    return [{"text": text.strip(), "page": None, "source": filepath.name}]


def parse_json(filepath: Path) -> list[dict]:
    raw = filepath.read_text(encoding="utf-8", errors="replace")
    try:
        data = _json.loads(raw)
    except Exception:
        return parse_text(filepath)
    text = _flatten_structured(data)
    if not text.strip():
        return []
    return [{"text": text.strip(), "page": None, "source": filepath.name}]


def parse_jsonl(filepath: Path) -> list[dict]:
    """Newline-delimited JSON. Each line becomes one section so the chunker
    can keep records semantically distinct."""
    results = []
    with open(filepath, "r", encoding="utf-8", errors="replace") as f:
        for i, line in enumerate(f, start=1):
            line = line.strip()
            if not line:
                continue
            try:
                obj = _json.loads(line)
                text = _flatten_structured(obj)
            except Exception:
                text = line
            if text.strip():
                results.append({"text": text.strip(), "page": i, "source": filepath.name})
    return results


def _flatten_structured(data, prefix: str = "") -> str:
    """Flatten a JSON/YAML-style nested structure into readable text.
    Preserves keys as labels so retrieval can match on field names too."""
    lines = []
    if isinstance(data, dict):
        for k, v in data.items():
            key_path = f"{prefix}.{k}" if prefix else str(k)
            if isinstance(v, (dict, list)):
                nested = _flatten_structured(v, key_path)
                if nested:
                    lines.append(nested)
            else:
                lines.append(f"{key_path}: {v}")
    elif isinstance(data, list):
        for i, item in enumerate(data):
            key_path = f"{prefix}[{i}]" if prefix else f"[{i}]"
            if isinstance(item, (dict, list)):
                nested = _flatten_structured(item, key_path)
                if nested:
                    lines.append(nested)
            else:
                lines.append(f"{key_path}: {item}")
    else:
        lines.append(f"{prefix}: {data}" if prefix else str(data))
    return "\n".join(lines)


# Code-language extensions get raw-text parsing for v1. The 300/50-word
# chunker handles them adequately. v2 will add AST-aware chunking via
# tree-sitter once chunking strategy is revisited across the product suite.
# List mirrors the graphify 5/11 file-type matrix (36+ code languages).
_CODE_EXTENSIONS = {
    # General-purpose
    ".py", ".pyi", ".ipynb",
    ".ts", ".tsx", ".js", ".jsx", ".mjs", ".cjs",
    ".go", ".rs", ".java", ".kt", ".scala",
    ".c", ".cpp", ".cc", ".cxx", ".h", ".hpp", ".hh",
    ".rb", ".cs", ".php", ".swift",
    ".lua", ".luau", ".zig",
    # Shell & scripting
    ".sh", ".bash", ".zsh", ".fish",
    ".ps1", ".psm1", ".bat", ".cmd",
    # Functional & ML
    ".ex", ".exs", ".jl", ".elm", ".clj", ".cljs", ".hs",
    # System & DSL
    ".sql", ".graphql", ".gql",
    ".m", ".mm", ".dart",
    # Web frameworks
    ".vue", ".svelte",
    # Build
    ".groovy", ".gradle",
    # Hardware / Verilog
    ".v", ".sv",
    # Fortran
    ".f", ".f90", ".f95", ".f03", ".f08",
    # Pascal / Delphi
    ".pas", ".pp", ".dpr", ".dpk", ".lpr", ".inc", ".dfm", ".lfm", ".lpk",
    # Infra-as-code
    ".tf", ".hcl",
    # GPU / shader
    ".cu", ".cuh", ".wgsl", ".glsl", ".hlsl",
}


def parse_file(filepath: Path) -> list[dict]:
    ext = filepath.suffix.lower()
    parsers = {
        # Existing
        ".pdf": parse_pdf,
        ".docx": parse_docx,
        ".doc": parse_docx,
        ".xlsx": parse_xlsx,
        ".xls": parse_xlsx,
        ".md": parse_markdown,
        ".mdx": parse_markdown,
        ".qmd": parse_markdown,
        ".txt": parse_text,
        ".rtf": parse_rtf,
        # New 2026-06-14 (graphify matrix)
        ".html": parse_html,
        ".htm": parse_html,
        ".pptx": parse_pptx,
        ".ppt": parse_pptx,
        ".yaml": parse_yaml,
        ".yml": parse_yaml,
        ".json": parse_json,
        ".jsonl": parse_jsonl,
        ".ndjson": parse_jsonl,
        ".rst": parse_text,  # raw-text v1; docutils-aware v2
    }
    # Code-extension dispatch
    if ext in _CODE_EXTENSIONS:
        return parse_text(filepath)
    parser = parsers.get(ext)
    if parser:
        return parser(filepath)
    # Fallback: try as text
    return [{"text": filepath.read_text(errors="replace"), "page": None, "source": filepath.name}]


def chunk_lines(text: str, chunk_size: int = 300, overlap: int = 50) -> list[str]:
    """Split text into word-budgeted chunks WITHOUT destroying its line structure.

    ⚠ THE ONE IMPLEMENTATION. There were three: this logic (correct, line-aware), and two
    copies in `cartridge_builder.chunk_text` and `forge.chunk_text` that did
    `" ".join(text.split())` and flattened every newline. Only the Cart Builder path had been
    fixed, so a cart's formatting depended on which script happened to build it -- the redwood
    demo carts came out of `build_office.py` -> `cartridge_builder.chunk_text` and arrived as
    a wall of text: 1709 of 1712 passages with zero newlines (2026-08-16).

    WHY IT MATTERS BEYOND TIDINESS. Markdown is LINE-BASED. Headings, bullets, and GFM tables
    all need line breaks to parse, and `PassageModal` renders with react-markdown +
    remark-gfm. Flattening leaves `#` and `|` sitting inline as literal punctuation, which is
    exactly the "jumbled together unformatted" the modal shows. Andy asked for indented
    bulleted lists on 2026-05-12; the display side has been able to do it since.

    ⚠ THE OLD VERSION LOOKED FINE ON SHORT INPUT. Under `chunk_size` it returned the text
    unmodified, so anything below the threshold kept its formatting and only long documents
    broke. That is why it read as an intermittent display problem rather than a chunker bug.

    Boundaries fall BETWEEN lines, never mid-line, so a table row or list item is never split.
    A single line longer than `chunk_size` is taken atomically -- one oversize chunk beats a
    mangled table row.
    """
    lines = text.split("\n")
    line_wc = [len(l.split()) for l in lines]
    n = len(lines)
    chunks: list[str] = []
    i = 0
    while i < n:
        j = i
        budget = 0
        while j < n and (budget == 0 or budget + line_wc[j] <= chunk_size):
            budget += line_wc[j]
            j += 1
        chunk = "\n".join(lines[i:j]).strip()
        if chunk:
            chunks.append(chunk)
        if j >= n:
            break
        # Walk backward from j to build an `overlap`-word tail that becomes the next chunk's
        # prefix. Guard: always advance at least one line so pathological inputs terminate.
        back = j
        overlap_wc = 0
        while back > i + 1 and overlap_wc < overlap:
            back -= 1
            overlap_wc += line_wc[back]
        i = back
    return chunks


def chunk_texts(sections: list[dict], chunk_size: int = 300, overlap: int = 50) -> list[dict]:
    """Split parsed sections into overlapping word-budgeted chunks.

    Line-aware: chunk boundaries fall between lines, never mid-line, so
    markdown structure (tables, lists, headings, paragraph breaks) survives
    intact and the passage viewer's react-markdown + remark-gfm renders
    them as-intended. Overlap is expressed as trailing lines whose combined
    word count is >= `overlap`.

    previous implementation (`text.split()` then
    `" ".join()`) destroyed all newlines, turning Docling OCR output into
    a single wall of piped text that remark-gfm parsed as paragraph, not
    table. This fix restores structural markdown.
    """
    chunks = []
    for section in sections:
        section_text = scrub_lone_surrogates(section["text"])
        if len(section_text.split()) <= chunk_size:
            chunks.append({**section, "text": section_text})
            continue
        for part, piece in enumerate(chunk_lines(section_text, chunk_size, overlap)):
            chunks.append({
                "text": piece,
                "page": section.get("page"),
                "source": section["source"],
                "part": part,
            })
    return chunks
